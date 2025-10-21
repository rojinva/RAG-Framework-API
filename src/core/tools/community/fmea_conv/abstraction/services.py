import os
import json
import pandas as pd
from src.models.tool import ToolArtifact
from langchain_core.retrievers import BaseRetriever
from src.core.retrievers.retriever import AzureAISearchRetriever
from src.core.tools.community.fmea_conv.abstraction.constant import StaticData, StaticFilters
from src.core.tools.community.fmea_conv.abstraction.fmea_agent_tool_spec_model import FMEAToolSpec, ToolInput, FMEARetrieverSpec
from src.core.base import LamBotTool
from io import BytesIO
from datetime import datetime
from langchain_openai import AzureChatOpenAI
from pptx import Presentation
from src.core.database.lambot import LamBotMongoDB
from src.models.constants import LanguageModelName
from src.clients.azure import openai_token_provider

llm_config_service = LamBotMongoDB.get_instance().language_model_config_db

class FmeaDataHandler: 
    @staticmethod
    def clean_and_rename_columns(df):
        """Clean and rename columns in the DataFrame."""
        df.columns = df.columns.str.replace("\n", "", regex=False).str.strip()

        if StaticFilters.RECOMMENDEDACTION in df.columns:
            df.rename(columns={StaticFilters.RECOMMENDEDACTION: StaticFilters.RECOMMENDEDCORRECTIVEACTION}, inplace=True)

        if StaticFilters.CORRECTIVEACTIONS in df.columns:
            df.rename(columns={StaticFilters.CORRECTIVEACTIONS: StaticFilters.RECOMMENDEDCORRECTIVEACTION}, inplace=True)
            df = df.iloc[1:]

        if StaticFilters.INTENDEDFUNCTION in df.columns:
            df.rename(columns={StaticFilters.INTENDEDFUNCTION: StaticFilters.ITEMFUNCTION}, inplace=True)

        return df
    
    @staticmethod
    def find_header_row(df, target_keywords):
        """Find the header row based on target keywords."""
        for i, row in df.iterrows():
            if any(keyword in row.values for keyword in target_keywords):
                skip_columns = next((index for index, value in enumerate(row) if value in target_keywords), 0)
                return i+1, skip_columns
        return None, 0

    @staticmethod
    def extract_fields_from_filled_fmea(file_bytes):
        df = pd.read_excel(BytesIO(file_bytes))
        
        if 'Type' in df.columns:
            return df
        else:    
            preview_df = pd.read_excel(BytesIO(file_bytes), sheet_name='FMEA', nrows=10)

            target_keywords = {
                StaticFilters.ITEMFUNCTION,
                StaticFilters.INTENDEDFUNCTION,
                "\n" + StaticFilters.ITEMFUNCTION,
                "\n" + StaticFilters.INTENDEDFUNCTION
            }

            header_row, skip_columns = FmeaDataHandler.find_header_row(preview_df, target_keywords)
            if header_row is None:
                return "Unable to process uploaded FMEA excel file, please make sure it has sheet named FMEA"

            df = pd.read_excel(BytesIO(file_bytes), sheet_name='FMEA', skiprows=header_row).iloc[:, skip_columns:skip_columns+13]
            df = FmeaDataHandler.clean_and_rename_columns(df)

            df = df.dropna(subset=[StaticFilters.ITEMFUNCTION])
        return df
    
    @staticmethod
    def create_design_requirements_from_pptx_bytes(pptx_bytes):
        try:
            pptx_stream = BytesIO(pptx_bytes)
            prs = Presentation(pptx_stream)
        except Exception as e:
            error_message = f"Failed to read uploaded PDR slide, please verify the file format and content: {str(e)}"
            return {"status": "error", "message": error_message}
            
        try:
            tables = []
            component_value = None
            part_value = None
            item_function_list = []
            boundary_conditions_list = []
            interaction_type_list = []
            pg_name = None
            bu_name = None
            tool = None

            for slide_number, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = []
                        table = shape.table

                        for row in table.rows:
                            table_data.append([cell.text.strip() for cell in row.cells])

                        df = pd.DataFrame(table_data)
                        df.attrs[StaticFilters.SLIDENUMBER] = slide_number
                        tables.append(df)
        except Exception as e:
            error_message = f"Failed to extract tables from uploaded PDR slides, please verify the file format and content: {str(e)}"
            return {"status": "error", "message": error_message}

        try:
            for function_item_df in tables:
                slide_num = function_item_df.attrs[StaticFilters.SLIDENUMBER]

                if slide_num == 1:
                    key_row = function_item_df[function_item_df.iloc[:, 0] == StaticFilters.COMPONENTKEY]
                    if not key_row.empty:
                        component_value = key_row.iloc[0, 1]
                        component_value = component_value.title()
                        
                    part_row = function_item_df[function_item_df.iloc[:, 0] == StaticFilters.PARTKEY]
                    if not part_row.empty:
                        part_value = part_row.iloc[0, 1]
                        part_value = part_value.strip()

                    pg_row = function_item_df[function_item_df.iloc[:, 0] == StaticFilters.PGNAME]
                    if not pg_row.empty:
                        pg_name = pg_row.iloc[0, 1]
                        pg_name = pg_name.strip()

                    bu_row = function_item_df[function_item_df.iloc[:, 0] == StaticFilters.BUNAME]
                    if not bu_row.empty:
                        bu_name = bu_row.iloc[0, 1]
                        bu_name = bu_name.strip()

                    tool_row = function_item_df[function_item_df.iloc[:, 0] == StaticFilters.TOOLNAME]
                    if not tool_row.empty:
                        tool = tool_row.iloc[0, 1]
                        tool = tool.strip()

                if 2 <= slide_num <= 6 and len(function_item_df.columns) >= 5:                    
                    function_item_df[StaticFilters.ITEMFUNCTION] = function_item_df.iloc[:, 1] + " " + function_item_df.iloc[:, 2]
                    item_function_list.extend(function_item_df[StaticFilters.ITEMFUNCTION].tolist()[1:])
                    
                    boundary_conditions_index = None
                    interaction_type_index = None
                    
                    for i, row in function_item_df.iterrows():
                        if "Boundary Conditions" in row.values:
                            boundary_conditions_index = 3
                        if "Interaction Type" in row.values:
                            interaction_type_index = 4

                    if boundary_conditions_index is not None:
                        boundary_conditions_list.extend(function_item_df.iloc[1:, boundary_conditions_index].tolist())
                    else:
                        boundary_conditions_list.extend([""] * (len(function_item_df) - 1))

                    if interaction_type_index is not None:
                        interaction_type_list.extend(function_item_df.iloc[1:, interaction_type_index].tolist())
                    else:
                        interaction_type_list.extend([""] * (len(function_item_df) - 1))
        except Exception as e:
            error_message = f"Failed to process tables and extract item functions, boundary conditions, and interaction types. Please verify the file format and content: {str(e)}"
            return {"status": "error", "message": error_message}

        if not component_value or not part_value or not item_function_list or not boundary_conditions_list or not interaction_type_list:
            error_message = "Failed to extract required data from uploaded PDR slides. Ensure the file contains valid tables and data."
            return {"status": "error", "message": error_message}

        return component_value, part_value, pg_name, bu_name, tool, item_function_list, boundary_conditions_list, interaction_type_list
    
    @staticmethod
    def combine_filled_and_generated_fmea(filled_fmea, df):
        try:
            if "Type" in filled_fmea.columns:
                filled_fmea["Type"] = filled_fmea["Type"].replace(['', None], "AI Filled")
                df["Type"] = "AI Filled-Updated"
            else:
                filled_fmea["Type"] = "Manually Filled"
                df["Type"] = "AI Filled"
                
            if "Date" not in filled_fmea.columns:
                filled_fmea["Date"] = "-"
            current_timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            df["Date"] = current_timestamp
            
            combined_fmea = pd.concat([filled_fmea, df], ignore_index=True)
            return combined_fmea
        except Exception as e:
            error_message = f"Failed to merge draft fmea with ai filled fmea: {str(e)}"
            return {"status": "error", "message": error_message}
        
    
    @staticmethod
    def configure_single_retriever(tool_spec: FMEARetrieverSpec) -> BaseRetriever:
        """Configure and return the AzureAISearchRetriever based on the provided tool configuration."""
        azure_search_config = tool_spec.search_config
        return AzureAISearchRetriever(
            name=tool_spec.index_name,
            index_name=tool_spec.index_name,
            tool_name=tool_spec.index_name,
            citation_field_mappings={},
            search_api_key=os.getenv("SEARCH_API_KEY"),
            search_api_base=os.getenv("SEARCH_API_BASE"),
            search_api_version=os.getenv("SEARCH_API_VERSION"),
            azure_search_config=azure_search_config,
            top_k=azure_search_config.get("top"),
            formatter=None,
            redact_pii = False
        )
    
    @staticmethod
    def configure_retrievers(tool_spec: FMEAToolSpec) -> BaseRetriever:
        """Configure and return each of the retrievers"""
        historical_fmea_retriever = FmeaDataHandler.configure_single_retriever(
            tool_spec=tool_spec.historical_fmea_retriever_spec
        )
        return [historical_fmea_retriever]
    
    @staticmethod
    def generate_output_with_llm_call(context, query=None):
        gpt5_chat_llm_config = llm_config_service.fetch_language_model(LanguageModelName.GPT_41)
        llm = AzureChatOpenAI(
            azure_ad_token_provider=openai_token_provider,
            azure_endpoint=gpt5_chat_llm_config.endpoint,
            api_version=gpt5_chat_llm_config.api_version,
            azure_deployment=gpt5_chat_llm_config.deployment_name,
            model=gpt5_chat_llm_config.name,
            temperature=0.0,
            seed=42,
            streaming=False,
        )

        if query:
            input_text = f"{context} {query}"
        else:
            input_text = context
        response = llm.invoke(input_text)
        return response
    
    @staticmethod
    def extract_filled_fmea_content(file_bytes):
        """Extract content from the Filled FMEA Excel file."""
        df = pd.read_excel(BytesIO(file_bytes))
        return df.to_string(index=False)
    
    @staticmethod
    def extract_pdr_slide_content(file_bytes):
        """Extract content from the Design Requirement PDR Slide."""
        pptx_stream = BytesIO(file_bytes)
        prs = Presentation(pptx_stream)
        slide_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_content.append(shape.text)

        return "\n".join(slide_content)
    
    @staticmethod
    def map_predefined_root_causes(predefined_root_causes, interaction_type_list):
        try:
            if not predefined_root_causes:
                return []
            
            if not interaction_type_list:
                return []
            
            root_causes_list = []
            for interaction in interaction_type_list:
                interaction_elements = interaction.split('||')
                
                combined_root_causes = []
                for element in interaction_elements:
                    element = element.strip()
                    if element in predefined_root_causes:
                        combined_root_causes.extend(predefined_root_causes[element])
                    else:
                        pass

                root_causes_list.append(", ".join(combined_root_causes))
            return root_causes_list
        except Exception as e:
            error_message = f"Unable to map pre-defined root causes: {str(e)}"
            return {"status": "error", "message": error_message}
    
    @staticmethod
    def parse_missing_line_items(missing_line_items):
        try:
            if missing_line_items.startswith("```json"):
                raw_missing_line_items  = missing_line_items[len("```json"):].strip()
            elif missing_line_items.startswith("```"):
                raw_missing_line_items  = missing_line_items[len("```"):].strip()
            else:
                raw_missing_line_items  = missing_line_items.strip()
            
            cleaned_missing_line_items = raw_missing_line_items .rstrip("```").strip()
            missing_line_items_list = json.loads(cleaned_missing_line_items)
            
            fmea_line_items = []
            for cause in missing_line_items_list:
                fmea_line_items.append({
                    StaticFilters.ITEMFUNCTION: cause.get(StaticFilters.ITEMFUNCTION, None),
                    StaticFilters.POTENTIALFAILUREMODE: cause.get(StaticFilters.POTENTIALFAILUREMODE, None),
                    StaticFilters.POTENTIALEFFECTOFFAILURE: cause.get(StaticFilters.POTENTIALEFFECTOFFAILURE, None),
                    StaticFilters.SEVERITY: cause.get(StaticFilters.SEVERITY, None),
                    StaticFilters.POTENTIALCUASEOFMECHANISM: cause.get(StaticFilters.POTENTIALCUASEOFMECHANISM, None),
                    StaticFilters.FAILURECATEGORY: cause.get(StaticFilters.FAILURECATEGORY, None),
                    StaticFilters.DESIGNPREVENTION: cause.get(StaticFilters.DESIGNPREVENTION, None),
                    StaticFilters.OCCURRENCE: cause.get(StaticFilters.OCCURRENCE, None),
                    StaticFilters.CLASS: cause.get(StaticFilters.CLASS, None),
                    StaticFilters.DESIGNCONTROL: cause.get(StaticFilters.DESIGNCONTROL, None),
                    StaticFilters.DETECTION: cause.get(StaticFilters.DETECTION, None),
                    StaticFilters.RPN: cause.get(StaticFilters.RPN, None),
                    StaticFilters.RECOMMENDEDCORRECTIVEACTION: cause.get(StaticFilters.RECOMMENDEDCORRECTIVEACTION, None)
                })
            return fmea_line_items
        except Exception as e:
            error_message = f"LLM output parsing failed: {str(e)}"
            return {"status": "error", "message": error_message}